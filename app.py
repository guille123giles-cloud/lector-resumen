import streamlit as st
import PyPDF2
import docx
import pptx
import tempfile
import asyncio
import edge_tts
import pytesseract
from PIL import Image
import textwrap
import os
import platform
import base64

# --- 0. CONFIGURACIÓN DE LA PÁGINA ---
st.set_page_config(page_title="Lector IA", page_icon="🎧", layout="wide")

# --- DEDICATORIA ---
st.markdown("<p style='text-align: center; font-size: 14px; font-style: italic; color: #e0c3fc;'>Para la más linda del mundo, Pili ❤️</p>", unsafe_allow_html=True)

# --- CONFIGURACIÓN DE TESSERACT MULTIPLATAFORMA ---
if platform.system() == "Windows":
    # Cambiá esta ruta si instalaste Tesseract en otro lugar en tu PC
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# --- 1. FUNCIONES PARA EXTRAER TEXTO ---
def extraer_texto_pdf(archivo):
    lector = PyPDF2.PdfReader(archivo)
    texto = ""
    for pagina in lector.pages:
        if pagina.extract_text():
            texto += pagina.extract_text() + "\n"
    return texto

def extraer_texto_word(archivo):
    doc = docx.Document(archivo)
    return "\n".join([parrafo.text for parrafo in doc.paragraphs])

def extraer_texto_pptx(archivo):
    presentacion = pptx.Presentation(archivo)
    texto = ""
    for diapositiva in presentacion.slides:
        for forma in diapositiva.shapes:
            if hasattr(forma, "text"):
                texto += forma.text + "\n"
        texto += "\n"
    return texto

def extraer_texto_imagen(archivo):
    imagen = Image.open(archivo)
    # En la nube usa 'spa' instalado vía packages.txt
    return pytesseract.image_to_string(imagen, lang='spa')

# --- FUNCIÓN ASÍNCRONA PARA GENERAR AUDIO ---
async def generar_audio_largo(texto, ruta_salida, velocidad_tts, voz_tts):
    fragmentos = textwrap.wrap(texto, width=3000, replace_whitespace=False)
    barra_progreso = st.progress(0)
    texto_estado = st.empty()
    archivos_temporales = []
    
    for i, fragmento in enumerate(fragmentos):
        texto_estado.text(f"Procesando parte {i+1} de {len(fragmentos)}...")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
            temp_path = fp.name
            archivos_temporales.append(temp_path)
        
        comunicador = edge_tts.Communicate(fragmento, voz_tts, rate=velocidad_tts)
        await comunicador.save(temp_path)
        barra_progreso.progress((i + 1) / len(fragmentos))
        
    with open(ruta_salida, 'wb') as archivo_final:
        for temp_path in archivos_temporales:
            with open(temp_path, 'rb') as f:
                archivo_final.write(f.read())
            os.remove(temp_path) 
    texto_estado.text("¡Audio listo!")

# --- 2. INTERFAZ VISUAL ---
st.title("🎧 Mi Lector de Resúmenes IA")
st.divider()

col_izq, col_der = st.columns([2, 1], gap="large")
texto_final = ""

with col_izq:
    st.subheader("📄 1. Cargá tu material")
    opcion = st.radio("¿Qué vas a usar?", ("Subir Archivo", "Pegar Texto"), horizontal=True)

    if opcion == "Subir Archivo":
        archivo_subido = st.file_uploader("PDF, Word, PPT o Foto", 
                                          type=["pdf", "docx", "pptx", "png", "jpg", "jpeg"])
        if archivo_subido:
            if archivo_subido.name.endswith(".pdf"):
                texto_final = extraer_texto_pdf(archivo_subido)
            elif archivo_subido.name.endswith(".docx"):
                texto_final = extraer_texto_word(archivo_subido)
            elif archivo_subido.name.endswith(".pptx"):
                texto_final = extraer_texto_pptx(archivo_subido)
            elif archivo_subido.name.endswith((".png", ".jpg", ".jpeg")):
                texto_final = extraer_texto_imagen(archivo_subido)
                st.image(archivo_subido, caption="Imagen cargada", use_container_width=True)
            
            st.success("¡Contenido leído!")
            with st.expander("Editar texto extraído"):
                texto_final = st.text_area("Texto:", value=texto_final, height=200)
    else:
        texto_final = st.text_area("Pegá tu resumen acá:", height=300)

with col_der:
    st.subheader("⚙️ 2. Ajustes")
    with st.expander("Ajustes de Voz y Velocidad", expanded=True):
        st.markdown("**🗣️ Voz**")
        voces = {
            "🇦🇷 Tomás (Arg)": "es-AR-TomasNeural",
            "🇦🇷 Elena (Arg)": "es-AR-ElenaNeural",
            "🇪🇸 Álvaro (Esp)": "es-ES-AlvaroNeural",
            "🇲🇽 Dalia (Mex)": "es-MX-DaliaNeural"
        }
        voz_elegida = voces[st.radio("Elegí acento:", list(voces.keys()))]
        
        st.divider()
        st.markdown("**⚡ Velocidad**")
        velocidades = {
            "Muy Lento (0.5x)": "-50%",
            "Lento (0.75x)": "-25%",
            "Normal (1x)": "+0%",
            "Rápido (1.25x)": "+25%",
            "Repaso (1.5x)": "+50%",
            "Turbo (2.0x)": "+100%"
        }
        velocidad_elegida = velocidades[st.radio("Velocidad:", list(velocidades.keys()), index=2)]

# --- 3. PROCESAMIENTO Y REPRODUCTOR ---
st.divider()
_, col_btn, _ = st.columns([1, 2, 1])

with col_btn:
    if texto_final and st.button("🔊 GENERAR AUDIO", use_container_width=True):
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
                archivo_audio = fp.name
            
            asyncio.run(generar_audio_largo(texto_final, archivo_audio, velocidad_elegida, voz_elegida))
            
            with open(archivo_audio, "rb") as f:
                audio_bytes = f.read()
                # Reproductor compatible con iPhone
                audio_b64 = base64.b64encode(audio_bytes).decode()
                st.markdown(f'<audio controls style="width: 100%;"><source src="data:audio/mpeg;base64,{audio_b64}" type="audio/mpeg"></audio>', unsafe_allow_html=True)
                
                st.download_button("⬇️ Descargar MP3", data=audio_bytes, file_name="resumen.mp3", mime="audio/mp3", use_container_width=True)
        except Exception as e:
            st.error(f"Error: {e}")
